"""Idea Validation orchestrator + per-stage runners."""
import json
from pathlib import Path
from datetime import datetime
from core.config import VAULT_PATH
from modules.agents.claude_cli import run_claude


TEMPLATES_DIR = VAULT_PATH / "99_System" / "Templates" / "idea_validation"
IDEAS_DIR = VAULT_PATH / "1_Projects" / "98_Ideen"
CRITERIA_FILE = VAULT_PATH / "99_System" / "Idea_Decision_Criteria.md"

# Item #78: headless `claude -p` can't answer permission prompts, so every tool
# a stage needs must be pre-allowed — otherwise WebSearch is silently denied and
# the stage degrades to caveated training-knowledge output (see the quarantined
# 2026-06-05 Parking market study).
STAGE_ALLOWED_TOOLS = [
    "Read", "Write", "Edit", "Bash", "WebSearch", "WebFetch", "Glob", "Grep",
]


STAGES = [
    ("02", "market_study", "02_Market_Study.md"),
    ("03", "competitive_landscape", "03_Competitive_Landscape.md"),
    ("04", "business_model", "04_Business_Model.md"),
    ("05", "financial_model", "05_Financial_Model.md"),
    ("06", "business_plan", "06_Business_Plan.md"),
    ("07", "pitch_deck", "07_Pitch_Deck.md"),
    ("08", "risk_register", "08_Risk_Register.md"),
    ("09", "capability_gap", "09_Capability_Gap.md"),
    ("10", "time_capital", "10_Time_Capital.md"),
    ("11", "scorecard", "11_Decision_Scorecard.md"),
    ("12", "premortem", "12_Pre_Mortem.md"),
    ("12b", "steelman_rebuttal", "12b_Steelman_Rebuttal.md"),
]


def idea_folder(idea_name: str) -> Path:
    safe = "".join(c if c.isalnum() or c in "_-" else "_" for c in idea_name)
    return IDEAS_DIR / safe


def create_idea(idea_name: str, brief_text: str) -> Path:
    folder = idea_folder(idea_name)
    folder.mkdir(parents=True, exist_ok=True)
    (folder / "_sources").mkdir(exist_ok=True)
    (folder / "_archive").mkdir(exist_ok=True)
    (folder / "01_Idea_Brief.md").write_text(f"# Idea Brief — {idea_name}\n\n{brief_text}\n", encoding="utf-8")
    (folder / "_overrides.md").write_text("# Overrides\n\nDavid's manual corrections / assumptions. Every stage prompt reads this and RESPECTS it.\n", encoding="utf-8")
    (folder / "MASTER.md").write_text(_master_template(idea_name), encoding="utf-8")
    _write_state(folder, {})
    _check_overlaps(idea_name, brief_text, folder)
    return folder


def _master_template(idea_name: str) -> str:
    return f"""# {idea_name}

Status: BRIEF (Stage 1 complete)
Created: {datetime.now().isoformat()}

## Quick links
- [Idea brief](01_Idea_Brief.md)
- [Sources](_sources/)
- [Overrides](_overrides.md)

## Progress
(populated by validation runs)

## Final recommendation
(set after Stage 11 + 12 complete)
"""


def _state_file(folder: Path) -> Path:
    return folder / "_state.json"


def _read_state(folder: Path) -> dict:
    p = _state_file(folder)
    return json.loads(p.read_text()) if p.exists() else {}


def _write_state(folder: Path, state: dict) -> None:
    _state_file(folder).write_text(json.dumps(state, indent=2, default=str), encoding="utf-8")


def _stage_template(stage_key: str) -> str:
    # Templates are named "<stage_key>_<slug>.md" (e.g. "02_market_study.md").
    # Resolve the slug from the STAGES table so the right file is loaded.
    slug = next((s for sk, s, _ in STAGES if sk == stage_key), None)
    name = f"{stage_key}_{slug}.md" if slug else f"{stage_key}.md"
    p = TEMPLATES_DIR / name
    return p.read_text(encoding="utf-8") if p.exists() else f"# Stage {stage_key}\n\n(Template missing — using default.)\n"


def _substitute(template: str, idea_name: str, stage_key: str) -> str:
    """Fill `{idea_name}` / `{STAGE_NAME}` placeholders in a stage template.

    Uses .replace(), NOT .format() — templates contain literal braces
    (EXCEL_SPEC blocks, the pitch deck's "[TBD — from {stage}]" instruction)
    that .format() would choke on.
    """
    slug = next((s for sk, s, _ in STAGES if sk == stage_key), None)
    stage_name = slug.replace("_", " ").title() if slug else f"Stage {stage_key}"
    return template.replace("{idea_name}", idea_name).replace("{STAGE_NAME}", stage_name)


def _require_brief(folder: Path) -> None:
    """Gate: refuse to run a stage unless a substantive idea brief exists.

    Without this, stages fire on empty context and Claude (correctly) refuses
    to fabricate — but the refusal text used to get written as the stage output.
    """
    brief = folder / "01_Idea_Brief.md"
    if not brief.exists():
        raise FileNotFoundError(
            f"01_Idea_Brief.md missing in {folder.name} — write the idea brief before running stages."
        )
    body = "\n".join(
        line for line in brief.read_text(encoding="utf-8").splitlines()
        if not line.lstrip().startswith("#")
    ).strip()
    if len(body) < 80:
        raise ValueError(
            f"01_Idea_Brief.md in {folder.name} is too thin ({len(body)} chars body, need 80+) — "
            "flesh out the brief before running stages."
        )


# Flag at ANY length — after _substitute these can never appear in a real output.
_HARD_REFUSAL_MARKERS = (
    "{idea_name}", "{stage_name}",   # echoed unsubstituted placeholders
    "unsubstituted",
    "i won't fabricate", "won't fabricate", "i can't fabricate", "cannot fabricate",
    "would just be fabrication",
)

# Flag only in short outputs (real stage docs are 1500+ words; refusals are brief).
_SOFT_REFUSAL_MARKERS = (
    "i can't help", "i cannot help",
    "what i won't do", "i won't do",
    "i need two things from you", "which idea do you",
    "missing `01_idea_brief", "missing 01_idea_brief",
    "stage template",
)


def _looks_like_refusal(output: str) -> bool:
    """Heuristic: detect a Claude refusal / clarification reply.

    Real stage outputs are long, structured Markdown (H2 sections, 1500+ words).
    Refusals are conversational replies addressed to David — usually short, but
    sometimes long enough to need the placeholder/fabrication hard markers.
    Guard before any write_text / state stamping / xlsx-docx-pptx generation.
    """
    text = output.strip()
    if not text:
        return True
    low = text.lower()
    if any(m in low for m in _HARD_REFUSAL_MARKERS):
        return True
    is_short = len(text) < 4000
    if is_short and any(m in low for m in _SOFT_REFUSAL_MARKERS):
        return True
    # No H2 structure at all in a short output → chat reply, not a stage doc.
    if is_short and "\n## " not in text and not text.startswith("## "):
        return True
    return False


# A real stage doc reaches its first heading almost immediately; a chat lead-in
# ("I have enough triangulated data. Writing the study now." + a `---` rule)
# is short. Anything longer than this before the first heading is treated as
# legitimate prose and left alone.
_PREAMBLE_MAX_CHARS = 400


def _strip_preamble(output: str) -> str:
    """Drop a short conversational lead-in before the first Markdown heading.

    ``claude -p`` sometimes prefixes a stage doc with a chat line addressed to
    David plus a horizontal rule (seen in the 2026-06-05 Parking market study).
    Strip anything before the first heading, but only when that lead-in is
    short and heading-free — a doc that legitimately opens with prose stays
    untouched.
    """
    text = output.lstrip()
    if text.startswith("#"):
        return text
    idx = text.find("\n#")
    if idx == -1 or idx > _PREAMBLE_MAX_CHARS:
        return text
    return text[idx + 1:].lstrip()


# Phrases a stage output uses to confess it worked WITHOUT live web research.
# Lowercase substrings; matched anywhere in the output (these caveats appear in
# the preamble or a methodology block, not necessarily the first lines).
_NO_WEBSEARCH_MARKERS = (
    "no web search",
    "web search was unavailable",
    "web search unavailable",
    "websearch` permission denied",
    "websearch permission denied",
    "training data only",
    "training knowledge",
    "medium-low confidence",
    "without access to current",
    "model cutoff",
)


def _stage_requires_research(stage_key: str) -> bool:
    """True if the stage template asks for web search / cited sources."""
    template = _stage_template(stage_key).lower()
    return any(
        marker in template
        for marker in ("web search", "websearch", "source", "citation", "cite")
    )


def _require_websearch_quality(output: str, stage_key: str) -> None:
    """Gate (Item #78): refuse training-knowledge-only output on research stages.

    When WebSearch is denied (or skipped), Claude produces a caveated
    reconstruction from training data — structurally a valid stage doc, so
    ``_looks_like_refusal`` passes it, but every figure is estimate-grade. If
    the output confesses to that AND the stage template demands web research /
    citations, raise instead of writing — same contract as the refusal guard:
    no write_text, no completed stamp, no xlsx/docx/pptx generation.
    """
    low = output.lower()
    hit = next((m for m in _NO_WEBSEARCH_MARKERS if m in low), None)
    if hit and _stage_requires_research(stage_key):
        raise RuntimeError(
            f"Stage {stage_key} output looks like training-knowledge-only work "
            f"(marker: '{hit}') but the stage requires web research. Nothing was "
            "written. Check that WebSearch is available (STAGE_ALLOWED_TOOLS / "
            "network), then re-run."
        )


def _all_prior_outputs(folder: Path, up_to_stage_key: str) -> str:
    """Concatenate all prior stage output files for inclusion in prompt context."""
    parts = []
    # Always include brief and overrides
    for f in ["01_Idea_Brief.md", "_overrides.md"]:
        p = folder / f
        if p.exists():
            parts.append(f"\n## {f}\n\n{p.read_text(encoding='utf-8')}")
    # Add prior numbered stages
    for stage_key, _, filename in STAGES:
        if stage_key == up_to_stage_key:
            break
        p = folder / filename
        if p.exists():
            parts.append(f"\n## {filename}\n\n{p.read_text(encoding='utf-8')}")
    return "\n---\n".join(parts)


def run_stage(idea_name: str, stage_key: str) -> Path:
    folder = idea_folder(idea_name)
    _require_brief(folder)
    template = _substitute(_stage_template(stage_key), idea_name, stage_key)
    prior = _all_prior_outputs(folder, stage_key)
    # Add criteria file if scorecard stage
    extra = ""
    if stage_key == "11":
        if CRITERIA_FILE.exists():
            extra = f"\n\n## Decision criteria (apply these)\n\n{CRITERIA_FILE.read_text(encoding='utf-8')}"

    prompt = f"""{template}

---

# REFERENCE CONTEXT FOR THIS STAGE

{prior}
{extra}

---

# YOUR TASK

Produce the stage output now, following the template's structure exactly. Use web search heavily. Cite every claim.
"""
    # Live-progress marker: lets the workspace UI show which stage is in
    # flight. Underscore-prefixed so stage lookups skip it; cleared in the
    # finally whether the stage succeeds or raises.
    state = _read_state(folder)
    state["_running"] = {"stage_key": stage_key, "started_at": datetime.now().isoformat()}
    _write_state(folder, state)
    try:
        output = run_claude(prompt, timeout=900, allowed_tools=STAGE_ALLOWED_TOOLS)

        if _looks_like_refusal(output):
            raise RuntimeError(
                f"Stage {stage_key} for '{idea_name}' returned a refusal/clarification reply, "
                f"not a stage output ({len(output.strip())} chars). Nothing was written. "
                "Check 01_Idea_Brief.md has real content, then re-run."
            )
        _require_websearch_quality(output, stage_key)
        output = _strip_preamble(output)

        # Find output filename from STAGES table
        output_filename = next((fn for sk, _, fn in STAGES if sk == stage_key), f"stage_{stage_key}.md")
        output_path = folder / output_filename

        # Archive prior version if exists
        if output_path.exists():
            archive_path = folder / "_archive" / f"{output_filename.replace('.md', '')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md"
            archive_path.parent.mkdir(exist_ok=True)
            archive_path.write_text(output_path.read_text(encoding="utf-8"), encoding="utf-8")

        output_path.write_text(output, encoding="utf-8")

        # Update state
        state = _read_state(folder)
        state[stage_key] = {"completed_at": datetime.now().isoformat(), "filename": output_filename}
        # Mark downstream stages as stale
        stage_keys_in_order = [sk for sk, _, _ in STAGES]
        if stage_key in stage_keys_in_order:
            idx = stage_keys_in_order.index(stage_key)
            for downstream in stage_keys_in_order[idx+1:]:
                if downstream in state:
                    state[downstream]["stale"] = True
        _write_state(folder, state)

        # Special handling for Excel/Word/PPT outputs
        if stage_key == "05":
            _generate_excel_from_spec(folder, output)
        elif stage_key == "06":
            _generate_docx_from_md(folder, output)
        elif stage_key == "07":
            _generate_pptx_from_spec(folder, output)

        return output_path
    finally:
        state = _read_state(folder)
        if state.pop("_running", None) is not None:
            _write_state(folder, state)


def run_all(idea_name: str) -> dict:
    """Run all stages sequentially. Returns dict of completed paths."""
    results = {}
    for stage_key, _, _ in STAGES:
        results[stage_key] = str(run_stage(idea_name, stage_key))
    _update_master_after_full_run(idea_name)
    return results


def run_remaining(idea_name: str) -> dict:
    """Run only the stages that are missing or stale, in pipeline order.

    Unlike :func:`run_all`, completed up-to-date stages are skipped — so a
    pipeline that died (or was fixed) mid-run resumes where it left off
    instead of re-paying for every prior Claude call.
    """
    folder = idea_folder(idea_name)
    results = {}
    for stage_key, _, filename in STAGES:
        state = _read_state(folder)  # re-read: run_stage flags downstream stale
        entry = state.get(stage_key)
        if entry and not entry.get("stale") and (folder / filename).exists():
            continue
        results[stage_key] = str(run_stage(idea_name, stage_key))
    _update_master_after_full_run(idea_name)
    return results


def _check_overlaps(idea_name: str, brief_text: str, folder: Path) -> None:
    """Look for overlapping market studies in other ideas; add references to MASTER.md."""
    # Stub for now — implement after multiple ideas exist
    pass


def _generate_excel_from_spec(folder: Path, financial_model_md: str) -> None:
    """Parse the ## EXCEL_SPEC block from financial model output and generate .xlsx."""
    # Extract spec block
    import re
    m = re.search(r"## EXCEL_SPEC\s*\n(.*?)(?=\n##|\Z)", financial_model_md, re.DOTALL)
    if not m:
        return
    spec = m.group(1)
    # Use openpyxl to generate. For v1, write the spec into a sheet structure.
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Financial Model"
    # Naive parse — write each line as a row
    for i, line in enumerate(spec.splitlines(), start=1):
        ws.cell(row=i, column=1, value=line)
    wb.save(folder / "05_Financial_Model.xlsx")


def _generate_docx_from_md(folder: Path, business_plan_md: str) -> None:
    """Convert markdown business plan to .docx using python-docx."""
    from docx import Document
    doc = Document()
    for line in business_plan_md.splitlines():
        if line.startswith("# "):
            doc.add_heading(line[2:], level=0)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=1)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=2)
        elif line.strip().startswith("- "):
            doc.add_paragraph(line.strip()[2:], style="List Bullet")
        elif line.strip():
            doc.add_paragraph(line)
    doc.save(folder / "06_Business_Plan.docx")


def _generate_pptx_from_spec(folder: Path, pitch_deck_md: str) -> None:
    """Parse ## PPTX_SPEC block; generate .pptx."""
    import re
    from pptx import Presentation
    from pptx.util import Inches, Pt
    m = re.search(r"## PPTX_SPEC\s*\n(.*?)(?=\n##|\Z)", pitch_deck_md, re.DOTALL)
    if not m:
        return
    spec = m.group(1)
    prs = Presentation()
    # Parse slides — each "### Slide N: Title" starts a new slide
    slides = re.split(r"###\s+Slide\s+\d+:?\s*", spec)
    for slide_block in slides[1:]:
        lines = [l for l in slide_block.splitlines() if l.strip()]
        if not lines:
            continue
        title = lines[0]
        bullets = [l.strip()[2:] for l in lines[1:] if l.strip().startswith("- ")]
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = bullets[0] if bullets else ""
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
    prs.save(folder / "07_Pitch_Deck.pptx")


def _update_master_after_full_run(idea_name: str) -> None:
    folder = idea_folder(idea_name)
    # Read scorecard recommendation if available; write to MASTER.md
    scorecard = folder / "11_Decision_Scorecard.md"
    if scorecard.exists():
        text = scorecard.read_text(encoding="utf-8")
        # Naive extraction — find a "Recommendation:" line
        import re
        m = re.search(r"Recommendation:\s*(.+)", text)
        recommendation = m.group(1) if m else "(see scorecard)"
        master = folder / "MASTER.md"
        if master.exists():
            content = master.read_text(encoding="utf-8")
            content = re.sub(r"## Final recommendation\n.*?(?=\n##|\Z)",
                            f"## Final recommendation\n{recommendation}\n",
                            content, flags=re.DOTALL)
            master.write_text(content, encoding="utf-8")


def kill_idea(idea_name: str, reason: str) -> None:
    """Move idea to _killed/ archive with reason.

    The folder move can briefly fail with a Windows ``PermissionError`` when
    OneDrive holds a sync lock on a just-written file (.xlsx/.docx/.pptx). Retry
    a few times before giving up.
    """
    import shutil
    import time

    folder = idea_folder(idea_name)
    if not folder.exists():
        return
    killed_dir = IDEAS_DIR / "_killed" / folder.name
    killed_dir.parent.mkdir(parents=True, exist_ok=True)
    last_exc: Exception | None = None
    for attempt in range(5):
        try:
            shutil.move(str(folder), str(killed_dir))
            last_exc = None
            break
        except PermissionError as exc:  # OneDrive sync lock — transient
            last_exc = exc
            time.sleep(1.5 * (attempt + 1))
    if last_exc is not None:
        raise last_exc
    (killed_dir / "_kill_reason.md").write_text(
        f"# Killed {datetime.now().isoformat()}\n\n{reason}\n", encoding="utf-8"
    )


def list_ideas() -> list[dict]:
    if not IDEAS_DIR.exists():
        return []
    out = []
    for d in IDEAS_DIR.iterdir():
        if not d.is_dir() or d.name.startswith("_"):
            continue
        state = _read_state(d)
        out.append({
            "name": d.name,
            "path": d,
            # Underscore keys (e.g. the _running marker) are not stages.
            "stages_complete": sum(1 for k in state if not k.startswith("_")),
            "state": state,
        })
    return sorted(out, key=lambda x: x["name"])


def list_killed() -> list[dict]:
    killed_dir = IDEAS_DIR / "_killed"
    if not killed_dir.exists():
        return []
    out = []
    for d in killed_dir.iterdir():
        if not d.is_dir():
            continue
        reason_file = d / "_kill_reason.md"
        reason = reason_file.read_text(encoding="utf-8") if reason_file.exists() else ""
        out.append({"name": d.name, "path": d, "reason": reason})
    return out
